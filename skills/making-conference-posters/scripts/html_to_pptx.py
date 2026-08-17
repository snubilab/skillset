#!/usr/bin/env python3
"""Convert a rendered poster HTML into an EDITABLE PowerPoint deck.

Not a screenshot on a slide: every box becomes a real shape and every text
block a real text frame, so the text can be edited in PowerPoint.

Method: render in headless Chrome, read each element's *resolved* geometry and
style out of the DOM, then re-emit them as PPTX shapes at the same physical
coordinates. Chrome has already done all the layout maths, so nothing is
re-derived from the CSS.

  S=<this-skill>/scripts   # wherever this skill is installed
  uv run --with python-pptx $S/html_to_pptx.py poster.html --sheet 84.1x118.9 \
      --root '#poster' --out poster-editable.pptx

Caveat worth knowing: PowerPoint re-wraps text with its own metrics, so a line
may break differently than in the PDF. The PDF stays the print master; this is
the file people edit. Install the poster's font on any machine that opens it.
"""
import argparse, html as _html, http.server, json, os, re, socketserver, subprocess, sys, threading

try:
    from pptx import Presentation
    from pptx.util import Cm, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    sys.exit("needs python-pptx:  uv run --with python-pptx html_to_pptx.py ...")

CHROME_CANDIDATES = [
    os.environ.get("CHROME_BIN", ""),
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/Applications/Chromium.app/Contents/MacOS/Chromium",
    "google-chrome", "google-chrome-stable", "chromium", "chromium-browser",
    "/snap/bin/chromium",
]

PROBE = r"""
const root=document.querySelector(__ROOT__);
if(!root){document.title='PPTXNOROOT';throw new Error('no root');}
const R0=root.getBoundingClientRect();
const CM=96/2.54;
const cm=v=>Math.round((v/CM)*10000)/10000;
const rgb=s=>{const m=(s||'').match(/-?[\d.]+/g); if(!m) return null;
  const a=m.length>3?parseFloat(m[3]):1; if(a<0.02) return null;
  return [Math.round(+m[0]),Math.round(+m[1]),Math.round(+m[2])];};
const hex=c=>c?c.map(x=>x.toString(16).padStart(2,'0')).join(''):null;

// clip-path polygon -> points relative to the border box.
// Chrome leaves calc() unresolved here, and calc(100% - 3.4cm) contains spaces,
// so split on commas/spaces only at paren depth 0, then evaluate each term.
function splitTop(str,comma){
  const out=[]; let depth=0,cur='';
  for(const ch of str){
    if(ch==='(') depth++;
    else if(ch===')') depth--;
    const hit = depth===0 && (comma ? ch===',' : /\s/.test(ch));
    if(hit){ if(cur.trim()) out.push(cur.trim()); cur=''; }
    else cur+=ch;
  }
  if(cur.trim()) out.push(cur.trim());
  return out;
}
function lenVal(tok,size){
  let t=tok.trim();
  if(t.startsWith('calc(')) t=t.slice(5,-1);
  const toks=t.match(/[+-]|[\d.]+%|[\d.]+px|[\d.]+cm|[\d.]+/g)||[];
  let total=0,sign=1;
  for(const tk of toks){
    if(tk==='+'){sign=1;continue;}
    if(tk==='-'){sign=-1;continue;}
    let v;
    if(tk.endsWith('%'))       v=parseFloat(tk)/100*size;
    else if(tk.endsWith('cm')) v=parseFloat(tk)*CM;
    else                       v=parseFloat(tk);
    if(!isFinite(v)) return NaN;
    total+=sign*v; sign=1;
  }
  return total;
}
function poly(cs,r){
  const cp=cs.clipPath||'';
  if(!cp.startsWith('polygon')) return null;
  const pts=[];
  for(const pt of splitTop(cp.slice(cp.indexOf('(')+1,cp.lastIndexOf(')')),true)){
    const xy=splitTop(pt,false);
    if(xy.length<2) return null;
    const x=lenVal(xy[0],r.width), y=lenVal(xy[1],r.height);
    if(!isFinite(x)||!isFinite(y)) return null;   // fall back to a plain rect
    pts.push([cm(x),cm(y)]);
  }
  return pts.length>2?pts:null;
}
// multi-stop linear-gradient -> sampled segments (PPTX has no n-stop axial fill)
function grad(cs){
  const bi=cs.backgroundImage||'';
  if(!bi.includes('linear-gradient')) return null;
  const stops=[...bi.matchAll(/rgba?\(([^)]+)\)\s*([\d.]+)%/g)]
    .map(m=>({c:rgb('rgb('+m[1]+')'),p:parseFloat(m[2])}));
  return stops.length>=2?stops:null;
}
const shapes=[],texts=[],images=[];
function walk(el,inText){
  const cs=getComputedStyle(el);
  if(cs.display==='none'||cs.visibility==='hidden') return;
  const r=el.getBoundingClientRect();
  const box={x:cm(r.left-R0.left),y:cm(r.top-R0.top),w:cm(r.width),h:cm(r.height)};
  if(el.tagName==='IMG'){
    images.push({...box,src:el.getAttribute('src')});
    return;
  }
  if(r.width>1&&r.height>1){
    const g=grad(cs);
    const fill=rgb(cs.backgroundColor);
    const bw=parseFloat(cs.borderTopWidth)||0;
    const bc=rgb(cs.borderTopColor);
    if(g)               shapes.push({...box,grad:g.map(s=>({c:hex(s.c),p:s.p}))});
    else if(fill||bw>0.5) shapes.push({...box,fill:hex(fill),
                                       line:bw>0.5?hex(bc):null,lw:cm(bw),
                                       poly:poly(cs,r)});
  }
  // a "text leaf" = an element with at least one non-empty direct text node
  const direct=[...el.childNodes].some(n=>n.nodeType===3&&n.textContent.trim());
  let captured=false;
  if(direct&&!inText){
    const runs=[];
    const push=(node,style,base)=>{
      const t=node.textContent.replace(/\s+/g,' ');
      if(!t.trim()&&!runs.length) return;
      runs.push({t,b:+style.fontWeight>=600,i:style.fontStyle==='italic',
                 sz:Math.round(parseFloat(style.fontSize)*0.75*10)/10,
                 c:hex(rgb(style.color)),f:style.fontFamily.split(',')[0].replace(/['"]/g,'').trim(),
                 sub:base==='sub',sup:base==='sup'});
    };
    (function collect(n,base){
      for(const ch of n.childNodes){
        if(ch.nodeType===3){ if(ch.textContent.trim()||runs.length) push(ch,getComputedStyle(n),base); }
        else if(ch.nodeType===1){
          const tag=ch.tagName.toLowerCase();
          collect(ch, tag==='sub'?'sub':tag==='sup'?'sup':base);
        }
      }
    })(el,null);
    if(runs.length){
      captured=true;
      const pl=parseFloat(cs.paddingLeft)||0, pr=parseFloat(cs.paddingRight)||0;
      const pt=parseFloat(cs.paddingTop)||0,  pb=parseFloat(cs.paddingBottom)||0;
      const bl=parseFloat(cs.borderLeftWidth)||0, br=parseFloat(cs.borderRightWidth)||0;
      texts.push({x:cm(r.left-R0.left+bl),y:cm(r.top-R0.top+pt),
                  w:cm(r.width-bl-br),h:cm(r.height-pt-pb),
                  align:cs.textAlign,lh:parseFloat(cs.lineHeight)/parseFloat(cs.fontSize)||1.2,
                  marL:cm(pl),indent:cm(parseFloat(cs.textIndent)||0),runs});
    }
  }
  for(const ch of el.children) walk(ch, inText||captured);
}
walk(root,false);
document.title='PPTX'+JSON.stringify({w:cm(R0.width),h:cm(R0.height),shapes,texts,images});
"""


def find_chrome():
    from shutil import which
    for c in CHROME_CANDIDATES:
        if not c:
            continue
        if os.path.exists(c):
            return c
        w = which(c)
        if w:
            return w
    sys.exit("Chrome/Chromium not found (set CHROME_BIN)")


def serve(directory):
    class H(http.server.SimpleHTTPRequestHandler):
        def __init__(self, *a, **k):
            super().__init__(*a, directory=directory, **k)
        def log_message(self, *a):
            pass
    srv = socketserver.TCPServer(("127.0.0.1", 0), H)
    threading.Thread(target=srv.serve_forever, daemon=True).start()
    return srv, srv.server_address[1]


def scrape(html_path, root_sel):
    d = os.path.dirname(os.path.abspath(html_path)) or "."
    js = PROBE.replace("__ROOT__", json.dumps(root_sel))
    src = open(html_path, encoding="utf-8").read()
    inject = ("<script>window.addEventListener('load',()=>document.fonts.ready"
              ".then(()=>setTimeout(()=>{" + js + "},400)));</script>")
    probe_html = (src.replace("</body>", inject + "</body>", 1)
                  if "</body>" in src else src + inject)
    wrapper = os.path.join(d, f"._pptx.{os.getpid()}.html")
    open(wrapper, "w", encoding="utf-8").write(probe_html)
    srv, port = serve(d)
    try:
        dom = subprocess.run([find_chrome(), "--headless", "--disable-gpu", "--no-sandbox",
                              "--virtual-time-budget=25000", "--dump-dom",
                              f"http://127.0.0.1:{port}/{os.path.basename(wrapper)}"],
                             capture_output=True, text=True, timeout=180).stdout
    finally:
        srv.shutdown()
        os.path.exists(wrapper) and os.remove(wrapper)
    if "<title>PPTXNOROOT</title>" in dom:
        sys.exit(f"--root {root_sel!r} matched no element")
    m = re.search(r"<title>PPTX(.*?)</title>", dom, re.S)
    if not m:
        sys.exit("probe did not run")
    return json.loads(_html.unescape(m.group(1))), d


ALIGN = {"left": PP_ALIGN.LEFT, "start": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
         "right": PP_ALIGN.RIGHT, "end": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}


def build(data, base_dir, out, sheet_w, sheet_h):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(sheet_w), Cm(sheet_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sh = slide.shapes

    for s in data["shapes"]:
        if s.get("grad"):
            # approximate an n-stop gradient with thin solid segments
            stops, N = s["grad"], 48
            for i in range(N):
                p = i / (N - 1) * 100
                lo = max([x for x in stops if x["p"] <= p], key=lambda x: x["p"], default=stops[0])
                hi = min([x for x in stops if x["p"] >= p], key=lambda x: x["p"], default=stops[-1])
                t = 0 if hi["p"] == lo["p"] else (p - lo["p"]) / (hi["p"] - lo["p"])
                col = [round(int(lo["c"][j*2:j*2+2], 16) * (1-t) + int(hi["c"][j*2:j*2+2], 16) * t)
                       for j in range(3)]
                seg = sh.add_shape(MSO_SHAPE.RECTANGLE, Cm(s["x"] + s["w"]*i/N), Cm(s["y"]),
                                   Cm(s["w"]/N + 0.02), Cm(s["h"]))
                seg.fill.solid(); seg.fill.fore_color.rgb = RGBColor(*col)
                seg.line.fill.background(); seg.shadow.inherit = False
            continue
        if s.get("poly") and all(p is not None and p[0] is not None and p[1] is not None for p in s["poly"]):
            ff = sh.build_freeform(Cm(s["x"] + s["poly"][0][0]), Cm(s["y"] + s["poly"][0][1]))
            ff.add_line_segments([(Cm(s["x"] + px), Cm(s["y"] + py)) for px, py in s["poly"][1:]],
                                 close=True)
            shp = ff.convert_to_shape()
        else:
            shp = sh.add_shape(MSO_SHAPE.RECTANGLE, Cm(s["x"]), Cm(s["y"]), Cm(s["w"]), Cm(s["h"]))
        if s.get("fill"):
            shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor.from_string(s["fill"].upper())
        else:
            shp.fill.background()
        if s.get("line"):
            shp.line.color.rgb = RGBColor.from_string(s["line"].upper())
            shp.line.width = Pt(max(0.75, s["lw"] * 28.35))
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False

    for im in data["images"]:
        p = os.path.join(base_dir, im["src"])
        if os.path.isfile(p):
            sh.add_picture(p, Cm(im["x"]), Cm(im["y"]), Cm(im["w"]), Cm(im["h"]))

    for t in data["texts"]:
        tb = sh.add_textbox(Cm(t["x"]), Cm(t["y"]), Cm(max(t["w"], 0.3)), Cm(max(t["h"], 0.3)))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        para = tf.paragraphs[0]
        para.alignment = ALIGN.get(t["align"], PP_ALIGN.LEFT)
        if t["runs"]:
            para.line_spacing = t["lh"]
        if t.get("marL") or t.get("indent"):
            # OOXML: marL is the left margin for every line, indent offsets the
            # first line. A negative indent equals a CSS hanging bullet.
            pPr = para._p.get_or_add_pPr()
            pPr.set("marL", str(int(Cm(max(0.0, t.get("marL", 0.0))))))
            pPr.set("indent", str(int(Cm(t.get("indent", 0.0)))))
        for r in t["runs"]:
            run = para.add_run()
            run.text = r["t"]
            f = run.font
            f.size = Pt(max(1, r["sz"])); f.bold = r["b"]; f.italic = r["i"]
            f.name = r["f"]
            if r.get("c"):
                f.color.rgb = RGBColor.from_string(r["c"].upper())
            if r.get("sub") or r.get("sup"):
                f._rPr.set("baseline", "-25000" if r["sub"] else "30000")

    prs.save(out)
    return len(data["shapes"]), len(data["texts"]), len(data["images"])


def main():
    p = argparse.ArgumentParser()
    p.add_argument("html")
    p.add_argument("--sheet", required=True, help="WxH in cm")
    p.add_argument("--root", default="body > *")
    p.add_argument("--out", required=True)
    a = p.parse_args()
    sw, sh_ = (float(x) for x in a.sheet.lower().split("x"))
    if max(sw, sh_) > 142.24:
        sys.exit(f"PowerPoint caps a slide at 142.24 cm; {max(sw, sh_)} cm will not fit")
    data, d = scrape(a.html, a.root)
    ns, nt, ni = build(data, d, a.out, sw, sh_)
    print(f"{a.out}\n  slide {sw} x {sh_} cm\n  {ns} shapes, {nt} editable text frames, {ni} images")


if __name__ == "__main__":
    main()
